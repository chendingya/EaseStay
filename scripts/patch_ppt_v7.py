from __future__ import annotations

from io import BytesIO
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.util import Pt
import xml.etree.ElementTree as ET


SRC = Path("docs/易宿酒店预订平台 - 三端架构设计与技术实现.pptx")
OUT = Path("docs/易宿酒店预订平台 - 三端架构设计与技术实现_修订版.pptx")


def _set_shape_text(slide, shape_id: int, text: str, font_pt: float | None = None) -> None:
    target = None
    for shp in slide.shapes:
        if shp.shape_id == shape_id:
            target = shp
            break
    if target is None or not getattr(target, "has_text_frame", False):
        return

    tf = target.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = "MiSans"
            if font_pt is not None:
                run.font.size = Pt(font_pt)


def _normalize_fonts_and_small_sizes(prs: Presentation) -> None:
    for sidx, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            if getattr(shp, "has_table", False):
                for row in shp.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                run.font.name = "MiSans"
                                if run.font.size is not None:
                                    pt = run.font.size.pt
                                    if 0 < pt < 9.5 and sidx not in (1, 2, 17):
                                        run.font.size = Pt(9.5)
                continue

            if not getattr(shp, "has_text_frame", False):
                continue
            for p in shp.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "MiSans"
                    if run.font.size is not None:
                        pt = run.font.size.pt
                        if 0 < pt < 9.5 and sidx not in (1, 2, 17):
                            run.font.size = Pt(9.5)


def _apply_content_patches(prs: Presentation) -> None:
    # Slide 1: avoid "答辩" wording.
    _set_shape_text(prs.slides[0], 20, "项目技术演示", 12)

    # Slide 4: stack details + route sample.
    _set_shape_text(prs.slides[3], 52, "React + Vite + Ant Design + react-router", 10.39)
    _set_shape_text(
        prs.slides[3],
        80,
        "router.use('/merchant/hotels', authRequired, requireRole('merchant'), merchantHotelRoutes)\n"
        "router.use('/admin/hotels', authRequired, requireRole('admin'), adminHotelRoutes)\n"
        "router.use('/requests', authRequired, requestRoutes)",
        8.9,
    )

    # Slide 6: align mobile architecture with markdown intent.
    _set_shape_text(prs.slides[5], 7, "四层纵向分层 + 横切能力", 12.71)
    _set_shape_text(prs.slides[5], 11, "App Shell", 12.71)
    _set_shape_text(prs.slides[5], 12, "根容器层", 9.5)
    _set_shape_text(prs.slides[5], 14, "职责: 全局用户状态 + 底部导航持久挂载", 10.17)
    _set_shape_text(prs.slides[5], 19, "Page Layer", 12.71)
    _set_shape_text(prs.slides[5], 20, "页面层", 9.5)
    _set_shape_text(prs.slides[5], 22, "职责: 11个页面独立路由,负责请求与业务回调", 10.17)
    _set_shape_text(prs.slides[5], 27, "Component Layer", 12.71)
    _set_shape_text(prs.slides[5], 28, "组件层", 9.5)
    _set_shape_text(prs.slides[5], 30, "职责: 布局组件 + ListContainer + Card纯展示", 10.17)
    _set_shape_text(prs.slides[5], 35, "Service Layer", 12.71)
    _set_shape_text(prs.slides[5], 36, "服务层", 9.5)
    _set_shape_text(prs.slides[5], 38, "职责: auth/favorites/request 统一封装", 10.17)
    _set_shape_text(prs.slides[5], 56, "横切能力", 11.44)
    _set_shape_text(prs.slides[5], 58, "全局鉴权: token校验,未登录跳转 /login", 9.5)
    _set_shape_text(prs.slides[5], 60, "列表抽象: createListByType 覆盖四类列表", 9.5)
    _set_shape_text(prs.slides[5], 62, "动效体系: 入场/骨架/下拉反馈统一封装", 9.5)
    _set_shape_text(prs.slides[5], 64, "复用收益: 页面只组装业务,不重复造轮子", 9.5)

    # Slide 8: code block readability.
    _set_shape_text(
        prs.slides[7],
        46,
        "const delay = animate ? `${Math.min(index, 10) * 20}ms` : '0ms'\n"
        "<View\n"
        "  className={`list-item${animate ? ' list-stagger-enter' : ''}`}\n"
        "  style={animate ? { animationDelay: delay } : undefined}\n"
        "/>",
        8.75,
    )

    # Slide 10: code block readability.
    _set_shape_text(
        prs.slides[9],
        52,
        "const DashboardBatchModals = lazy(() => import('../components/DashboardBatchModals.jsx'))\n"
        "const AuditTable = lazy(() => import('../components/audit/AuditTable.jsx'))\n"
        "const ReactECharts = lazy(() => import('echarts-for-react'))\n"
        "scheduleIdleTask(async () => {\n"
        "  setUnreadCount(await getUnreadCount())\n"
        "}, { timeout: 1500, fallbackDelay: 400 })",
        8.8,
    )

    # Slide 11: code block readability.
    _set_shape_text(
        prs.slides[10],
        58,
        "<RoomsTabBase roomTypes={rooms} i18nPrefix=\"merchant.room\"\n"
        "  showStatusColumn={false} showActionColumn={false} />\n"
        "<RoomsTabBase roomTypes={rooms} i18nPrefix=\"admin.room\"\n"
        "  showStatusColumn={true} showActionColumn={true}\n"
        "  onOpenDiscount={handleOpenDiscount}\n"
        "  onCancelDiscount={handleCancelDiscount} />",
        8.8,
    )

    # Slide 12: code block readability.
    _set_shape_text(
        prs.slides[11],
        48,
        "const authRequired = (req, res, next) => {\n"
        "  const token = req.headers.authorization?.split(' ')[1]\n"
        "  if (!token) return res.status(401).json({ error: 'Unauthorized' })\n"
        "  try { req.user = verifyToken(token); next() } catch { res.status(401).json({ error: 'Invalid token' }) }\n"
        "}\n"
        "const requireRole = (role) => (req, res, next) => req.user.role === role ? next() : res.status(403).json({ error: 'Forbidden' })",
        8.6,
    )

    # Slide 13: route sample complete.
    _set_shape_text(
        prs.slides[12],
        58,
        "router.use('/merchant/hotels', authRequired, requireRole('merchant'), merchantHotelRoutes)\n"
        "router.use('/admin/hotels', authRequired, requireRole('admin'), adminHotelRoutes)\n"
        "router.use('/requests', authRequired, requestRoutes)",
        8.9,
    )

    # Slide 14: clarify duplicate state labels and code block.
    _set_shape_text(prs.slides[13], 20, "取消(超时)", 9.5)
    _set_shape_text(prs.slides[13], 25, "取消(用户)", 9.5)
    _set_shape_text(prs.slides[13], 44, "重新提交", 9.5)
    _set_shape_text(prs.slides[13], 47, "待审核(复审)", 9.5)
    _set_shape_text(
        prs.slides[13],
        70,
        "const effectivePrice = calculateRoomPrice(roomType, { checkIn, checkOut })\n"
        "const total_price = roundToTwo(effectivePrice * quantity * nights)\n"
        "order.price_per_night = effectivePrice",
        8.8,
    )

    # Slide 15: remove off-topic section and improve code readability.
    _set_shape_text(prs.slides[14], 58, "关键收益", 11.87)
    _set_shape_text(prs.slides[14], 60, "实时库存: 无预扣减仍可准确返回可售量", 10.0)
    _set_shape_text(prs.slides[14], 62, "搜索排序: 城市标准化 + 相关性评分提升命中", 10.0)
    _set_shape_text(prs.slides[14], 64, "兼容策略: 全量/分页双模输出,旧调用零改造", 10.0)
    _set_shape_text(
        prs.slides[14],
        40,
        "query.lt('check_in', normalizedCheckOut)\n"
        "  .gt('check_out', normalizedCheckIn)\n"
        ";(data || []).forEach(row => {\n"
        "  const prev = map.get(row.room_type_id) || 0\n"
        "  map.set(row.room_type_id, prev + row.quantity)\n"
        "})",
        8.8,
    )


def _patch_arrowheads(path: Path, slide_indexes: list[int], color_hex: str = "39A9FF") -> None:
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    ET.register_namespace("a", ns["a"])
    ET.register_namespace("p", "http://schemas.openxmlformats.org/presentationml/2006/main")
    ET.register_namespace("r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships")

    with ZipFile(path, "r") as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    for idx in slide_indexes:
        slide_name = f"ppt/slides/slide{idx}.xml"
        if slide_name not in files:
            continue
        root = ET.fromstring(files[slide_name])
        updated = False
        for ln in root.findall(".//a:ln", ns):
            fill = ln.find("a:solidFill", ns)
            if fill is None:
                continue
            srgb = fill.find("a:srgbClr", ns)
            if srgb is None or srgb.get("val", "").upper() != color_hex:
                continue
            head = ln.find("a:headEnd", ns)
            if head is None:
                head = ET.SubElement(ln, f"{{{ns['a']}}}headEnd")
            head.set("type", "triangle")
            head.set("w", "med")
            head.set("len", "med")
            updated = True
        if updated:
            files[slide_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    buf = BytesIO()
    with ZipFile(buf, "w", ZIP_DEFLATED) as zout:
        for name, content in files.items():
            zout.writestr(name, content)

    path.write_bytes(buf.getvalue())


def main() -> None:
    prs = Presentation(SRC)
    _apply_content_patches(prs)
    _normalize_fonts_and_small_sizes(prs)
    prs.save(OUT)
    _patch_arrowheads(OUT, [4, 5, 9, 13])
    print(f"patched: {OUT}")


if __name__ == "__main__":
    main()
