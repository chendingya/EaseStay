from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn


def set_font_for_text_frame(text_frame, font_name: str) -> int:
    changed = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            rPr = run._r.get_or_add_rPr()
            before = run.font.name
            run.font.name = font_name

            # Ensure Office uses the same family for Latin / East Asian / Complex script.
            for tag in ("latin", "ea", "cs"):
                node = rPr.find(qn(f"a:{tag}"))
                if node is None:
                    node = OxmlElement(f"a:{tag}")
                    rPr.append(node)
                node.set("typeface", font_name)

            if before != font_name:
                changed += 1
    return changed


def set_font_in_pptx(path: Path, font_name: str) -> int:
    prs = Presentation(path)
    changed = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                changed += set_font_for_text_frame(shape.text_frame, font_name)
                continue

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        changed += set_font_for_text_frame(cell.text_frame, font_name)

    prs.save(path)
    return changed


def main() -> None:
    parser = argparse.ArgumentParser(description="Set global font family for a PPTX file.")
    parser.add_argument("pptx", type=Path, help="Target pptx path (in-place update)")
    parser.add_argument("--font", default="微软雅黑", help="Font family name")
    args = parser.parse_args()

    changed = set_font_in_pptx(args.pptx, args.font)
    print(f"updated_runs={changed}")


if __name__ == "__main__":
    main()
